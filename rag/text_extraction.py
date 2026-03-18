import os
import re
import sys
import argparse
import pdfplumber
import easyocr
import trafilatura
import io
import numpy as np
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter
from PIL import Image

current_script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(current_script_dir)
if project_root not in sys.path:
    sys.path.append(project_root)
try:
    from config import Config
except ImportError:
    print("錯誤: 找不到 config.py，請確認檔案結構。")
    sys.exit(1)

# 設定路徑
pdf_dirs = Config.PDF_DIR
test_output_dir = os.path.join(current_script_dir, "result")
reader = easyocr.Reader(['en', 'ch_tra'], gpu=False)

def perform_ocr(image_bytes):
    """執行 OCR 並回傳清理後的文字"""
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        img_np = np.array(img)
        # detail=0 只回傳純文字列表
        results = reader.readtext(img_np, detail=0)
        if results:
            text = " ".join(results).strip()
            return f" [圖片 OCR內容: {text}] "
    except Exception as e:
        print(f"      OCR 處理失敗: {e}")
    return ""

def get_header_footer_blacklist(pdf_path, threshold=5):
    """建立過濾黑名單"""
    candidates = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                h, w = page.height, page.width
                # 掃描上下各 10%
                for bbox in [(0, 0, w, h * 0.1), (0, h * 0.9, w, h)]:
                    crop = page.within_bbox(bbox)
                    text = crop.extract_text()
                    if text:
                        lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 2]
                        candidates.extend(lines)
    except Exception as e:
        print(f"黑名單掃描失敗: {e}")
    return {text for text, count in Counter(candidates).items() if count >= threshold}

def clean_text(text, blacklist=None):
    """執行正則清理與黑名單過濾"""
    if not text:
        return ""

    # 依照黑名單過濾
    if blacklist:
        lines = text.split('\n')
        # 如果該行文字存在於黑名單中，就移除
        lines = [line for line in lines if line.strip() not in blacklist]
        text = '\n'.join(lines)
    
    return text.strip()

def extract_docx_with_ocr(file_path):
    """解析 Word 並提取內嵌圖片進行 OCR"""
    doc = Document(file_path)
    full_text = []
    
    # 1. 處理段落文字
    for para in doc.paragraphs:
        full_text.append(para.text)
    
    # 2. 處理內嵌圖片 (Inline Shapes)
    for shape in doc.inline_shapes:
        try:
            # 取得圖片的二進位資料
            image_bytes = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            image_part = doc.part.related_parts[image_bytes]
            ocr_result = perform_ocr(image_part.blob)
            if ocr_result:
                full_text.append(ocr_result)
        except Exception:
            continue
            
    return "\n".join(full_text)

def extract_pptx_with_ocr(file_path):
    """解析 PPT 並依照投影片順序執行 OCR"""
    prs = Presentation(file_path)
    full_text = []
    
    for i, slide in enumerate(prs.slides):
        full_text.append(f"\n[Slide {i+1}]")
        # 依照座標排序 (由上而下)
        shapes = sorted(slide.shapes, key=lambda s: (s.top if hasattr(s, 'top') else 0))
        
        for shape in shapes:
            # 文字框處理
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text.strip())
            
            # 圖片 OCR 處理
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"    - Slide {i+1}: 偵測到圖片，正在 OCR...")
                ocr_result = perform_ocr(shape.image.blob)
                if ocr_result:
                    full_text.append(ocr_result)
                    
        # 備忘錄
        '''
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes: full_text.append(f"[Note]: {notes}")
        '''
            
    return "\n".join(full_text)

IMAGE_EXTS = (".jpg", ".jpeg", ".png")
def extract_image_only(file_path):
    """處理純圖片檔案的解析"""
    print(f"    - 執行純圖片 OCR: {os.path.basename(file_path)}")
    try:
        with open(file_path, "rb") as f:
            image_bytes = f.read()
            # 直接使用我們先前定義好的 perform_ocr 函式
            ocr_result = perform_ocr(image_bytes)
            # 移除 [圖片 OCR內容: ] 標籤，因為這整個檔案就是圖片
            # 或者是保留它以維持格式統一，這裡我選擇簡化它
            clean_result = ocr_result.replace("[圖片 OCR內容: ", "").replace("]", "").strip()
            return f"[Image File: {os.path.basename(file_path)}]\n{clean_result}"
    except Exception as e:
        print(f"  純圖片讀取失敗: {e}")
        return ""

def extract_pdf_with_ocr(file_path, skip_first_page=False):
    """解析 PDF：結合原生文字提取與圖片區域的 OCR"""
    # 取得該 PDF 的頁首頁尾黑名單
    blacklist = get_header_footer_blacklist(file_path)
    if blacklist:
        print(f"    [黑名單命中]: {blacklist}")

    pages_text = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                if skip_first_page and i == 0:
                    continue
                
                # A. 提取原生文字層
                raw_text = page.extract_text() or ""
                cleaned_raw = clean_text(raw_text, blacklist)
                
                # B. 處理頁面中的圖片物件
                ocr_results = []
                # page.images 包含了該頁所有嵌入圖片的座標
                if page.images:
                    print(f"    - Page {i+1}: 偵測到 {len(page.images)} 張圖片，執行 OCR...")
                    for img_obj in page.images:
                        try:
                            # 設定截圖邊界 (x0, top, x1, bottom)
                            bbox = (img_obj['x0'], img_obj['top'], img_obj['x1'], img_obj['bottom'])
                            
                            # 過濾掉太小的裝飾性圖標 (寬高小於 30 像素)
                            if (bbox[2] - bbox[0] < 30) or (bbox[3] - bbox[1] < 30):
                                continue
                                
                            # 使用 pdfplumber 裁切該區域並轉為影像
                            cropped_page = page.within_bbox(bbox).to_image(resolution=200)
                            
                            # 將影像轉為 bytes 傳給 OCR 處理
                            with io.BytesIO() as img_buffer:
                                cropped_page.save(img_buffer, format="PNG")
                                ocr_text = perform_ocr(img_buffer.getvalue())
                                if ocr_text:
                                    ocr_results.append(ocr_text)
                        except Exception as e:
                            # 避免單張圖片出錯中斷整頁解析
                            continue

                # C. 組合文字層與 OCR 層
                combined_content = f"[Page {i+1}]\n{cleaned_raw}"
                if ocr_results:
                    combined_content += "\n" + "\n".join(ocr_results)
                
                pages_text.append(combined_content)
                
    except Exception as e:
        print(f"  PDF OCR 解析失敗 {file_path}: {e}")
        
    return "\n\n".join(pages_text)

def extract_url(url):
    """從網址提取乾淨的內文"""
    print(f"    - 正在抓取網頁內容: {url}")
    try:
        # 下載網頁
        downloaded = trafilatura.fetch_url(url)
        if downloaded is None:
            return f"錯誤: 無法下載該網址內容 ({url})"
        
        # 提取內文 (已自動去除廣告、側邊欄)
        # include_comments=False 可以確保內容更乾淨
        result = trafilatura.extract(downloaded, include_comments=False, include_tables=True)
        
        if result:
            return f"[Source URL: {url}]\n{result.strip()}"
        else:
            return f"警告: 網址解析後無內容 ({url})"
            
    except Exception as e:
        print(f"  網頁抓取失敗: {e}")
        return ""

def extract_all_text(source, skip_first_page=False):
    """通用解析器分支"""
    # 判斷是否為網址
    if source.startswith(("http://", "https://")):
        return extract_url(source)
    
    # 檔案路徑判斷
    ext = os.path.splitext(source)[1].lower()
    
    if ext == ".pdf":
        return extract_pdf_with_ocr(source, skip_first_page)
    elif ext == ".docx":
        return extract_docx_with_ocr(source)
    elif ext == ".pptx":
        return extract_pptx_with_ocr(source)
    elif ext in IMAGE_EXTS:
        return extract_image_only(source)
    
    return ""

def save_text(text, out_dir, filename):
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, filename)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"  [已儲存]: {out_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="多格式文件與網頁轉純文字工具")
    parser.add_argument("--source", type=str, help="指定來源 (檔案路徑或 URL)")
    parser.add_argument("--category", type=str, help="指定特定分類資料夾 (例如: Computer)")
    args = parser.parse_args()

    output_root = Config.TXT_DIR
    supported_exts = (".pdf", ".docx", ".pptx", ".jpg", ".jpeg", ".png", ".bmp", ".webp")

    # 1. 判斷是否有指定特定來源 (--source)
    if args.source:
        source = args.source
        out_dir = os.path.join(output_root, args.category) if args.category else output_root
        
        # 模式 A: 處理網址 URL
        if source.startswith(("http://", "https://")):
            print(f"正在抓取網頁內容: {source}")
            result_text = extract_all_text(source)
            # URL 的存檔名稱處理（移除特殊字元）
            safe_name = re.sub(r'[^\w\s-]', '_', source.split("//")[-1])[:50] + ".txt"
            save_text(result_text, out_dir, safe_name)
            
        # 模式 B: 處理單一本地檔案
        elif os.path.isfile(source):
            ext = os.path.splitext(source)[1].lower()
            if ext in supported_exts:
                print(f"正在處理單一檔案: {source}")
                skip = any(cat in source for cat in ["Computer", "Physics"])
                result_text = extract_all_text(source, skip_first_page=skip)
                save_name = os.path.splitext(os.path.basename(source))[0] + ".txt"
                save_text(result_text, out_dir, save_name)
            else:
                print(f"錯誤: 不支援的副檔名 {ext}")
        
        # 模式 C: 處理指定目錄
        elif os.path.isdir(source):
            print(f"正在處理目錄: {source}")
            # 這裡可以調用下方的批次處理邏輯，或者直接遍歷該目錄
            for file in os.listdir(source):
                if file.lower().endswith(supported_exts):
                    p_path = os.path.join(source, file)
                    text = extract_all_text(p_path)
                    save_text(text, out_dir, os.path.splitext(file)[0] + ".txt")
        else:
            print(f"錯誤: 找不到來源 {source}")

    # 2. 預設模式：批次處理 Config 裡的 pdf_dirs
    else:
        print("未指定 --source，執行 Config 預設批次處理...")
        for pdf_dir in pdf_dirs:
            category = os.path.basename(os.path.normpath(pdf_dir))
            if args.category and args.category != category:
                continue
                
            print(f"\n--- 正在處理分類: {category} ---")
            out_dir = os.path.join(test_output_dir, category)
            skip_first = category in ["Computer", "Physics"]

            if not os.path.exists(pdf_dir):
                print(f"警告: 路徑不存在 {pdf_dir}")
                continue

            for file in os.listdir(pdf_dir):
                if file.lower().endswith(supported_exts):
                    p_path = os.path.join(pdf_dir, file)
                    print(f"處理中: {file}")
                    text = extract_all_text(p_path, skip_first_page=skip_first)
                    filename_only = os.path.splitext(file)[0]
                    save_text(text, out_dir, f"{filename_only}.txt")